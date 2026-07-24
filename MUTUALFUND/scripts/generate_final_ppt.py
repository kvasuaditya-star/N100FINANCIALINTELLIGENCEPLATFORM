import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Setup paths
BASE_DIR = Path(__file__).resolve().parent.parent
REPORTS_DIR = BASE_DIR / "reports"

# Helpers
def add_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=(255,255,255), alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=16, bold=False, color=(220,220,220),
             font_name="Calibri", space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name = font_name
    p.space_before = Pt(space_before)
    return p

def add_shape_bar(slide, left, top, width, height, r, g, b):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()

# Colors
BG       = (13, 14, 21)      # deep dark blue-grey
ACCENT   = (0, 180, 216)     # cyan
ACCENT2  = (114, 9, 183)     # purple
WHITE    = (255, 255, 255)
LIGHT    = (200, 200, 220)
GOLD     = (255, 196, 0)
GREEN    = (0, 200, 120)

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # SLIDE 1: Title
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0, 0, 13.333, 0.15, *ACCENT)
    add_shape_bar(sl, 0, 7.35, 13.333, 0.15, *ACCENT2)
    
    add_textbox(sl, 1, 2.0, 11, 1.2, "BLUESTOCK MUTUAL FUND ANALYTICS", 42, True, ACCENT, PP_ALIGN.CENTER)
    add_textbox(sl, 1, 3.2, 11, 0.6, "Comprehensive Portfolio Management & Financial Engineering", 22, False, WHITE, PP_ALIGN.CENTER)
    add_textbox(sl, 1, 4.2, 11, 0.4, "Capstone Project - Final Presentation", 18, True, GOLD, PP_ALIGN.CENTER)
    add_textbox(sl, 1, 5.8, 11, 0.5, "Author: Vasu | OS: Windows | Repository: MUTUALFUNDSDATAANALYTICS", 14, False, LIGHT, PP_ALIGN.CENTER)
    
    # -------------------------------------------------------------
    # SLIDE 2: Objectives & Context
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0.5, 0.5, 0.08, 0.6, *ACCENT)
    add_textbox(sl, 0.8, 0.4, 11, 0.7, "Project Objectives & Scope", 30, True, WHITE)
    
    tf = add_textbox(sl, 0.8, 1.5, 11.5, 5.0, "Core Business Goals", 20, True, ACCENT)
    add_para(tf, "• ETL Data Pipeline (D1/B1): Automated ingestion from live APIs (mfapi.in) and cleaning raw files.", 16, False, LIGHT)
    add_para(tf, "• Database Storage (D2): Load normalized tables into a relational SQLite schema for structured analytical querying.", 16, False, LIGHT)
    add_para(tf, "• Performance Analytics (D4): Mathematical formulation of CAGR, Sharpe ratio, Sortino, and OLS regression Alpha/Beta.", 16, False, LIGHT)
    add_para(tf, "• Advanced Financial Engineering (D6/B3/B4): Value at Risk (VaR), Conditional VaR, Monte Carlo projections, and Markowitz portfolio optimization.", 16, False, LIGHT)
    add_para(tf, "• Interactive Web Dashboard (D5/B2): Multi-page Streamlit web app loaded with interactive slicers, graphs, and recommender engines.", 16, False, LIGHT)

    # -------------------------------------------------------------
    # SLIDE 3: Database & ETL Architecture
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0.5, 0.5, 0.08, 0.6, *ACCENT)
    add_textbox(sl, 0.8, 0.4, 11, 0.7, "Database & ETL Pipeline Architecture", 30, True, WHITE)
    
    tf = add_textbox(sl, 0.8, 1.5, 5.5, 5.0, "Database Schema", 20, True, ACCENT)
    add_para(tf, "• dim_fund (Primary Key: amfi_code) - Scheme name, house, Category, expense ratio, launch date, risk classification.", 16, False, LIGHT)
    add_para(tf, "• fact_nav (PK: amfi_code, date) - Daily NAV time-series.", 16, False, LIGHT)
    add_para(tf, "• fact_transactions - Investor transactions logs (SIP, Lumpsum, Redemption, state, gender, income, KYC).", 16, False, LIGHT)
    add_para(tf, "• fact_performance - Trailing CAGR, Sharpe, Sortino, Alpha, Beta, Volatility, Max Drawdown.", 16, False, LIGHT)
    add_para(tf, "• portfolio_holdings - Assets and sectors weights for HHI calculation.", 16, False, LIGHT)
    
    tf2 = add_textbox(sl, 6.8, 1.5, 5.5, 5.0, "ETL Pipeline Automation", 20, True, GOLD)
    add_para(tf2, "• Live Queries: Queries mfapi.in API and fetches newest NAVs.", 16, False, LIGHT)
    add_para(tf2, "• Path Resolution: Replaces hardcoded strings with dynamic pathlib Pathing.", 16, False, LIGHT)
    add_para(tf2, "• Auto Schema: Recreates the SQLite db, builds tables, constraints, indexes, and loads data.", 16, False, LIGHT)
    add_para(tf2, "• Scheduling (B1): Background weekday daemon configured to execute every day at 8 PM.", 16, False, LIGHT)

    # -------------------------------------------------------------
    # SLIDE 4: Reindexing & Weekend NAV Handling
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0.5, 0.5, 0.08, 0.6, *ACCENT)
    add_textbox(sl, 0.8, 0.4, 11, 0.7, "NAV Reindexing & Weekend Handling", 30, True, WHITE)
    
    tf = add_textbox(sl, 0.8, 1.5, 11.5, 5.0, "Correct Holiday/Weekend Management", 20, True, ACCENT)
    add_para(tf, "• Problem: Mutual funds do not publish NAVs on weekends and public market holidays, causing time gaps in daily return metrics.", 16, False, LIGHT)
    add_para(tf, "• Solution: Reindex the NAV series to a continuous daily date index (calendar range including Saturday and Sunday).", 16, False, LIGHT)
    add_para(tf, "• Forward-Fill: Call pandas ffill() to copy the Friday NAV to Saturday and Sunday. This avoids return spikes and matches calendar realities.", 16, False, LIGHT)
    add_para(tf, "• CAGR Annualisation: Using calendar days is avoided. The CAGR calculation is annualized using trading days:", 16, False, LIGHT)
    add_para(tf, "  CAGR = (Ending NAV / Beginning NAV) ^ (252 / n_trading_days) - 1", 18, True, GOLD)
    add_para(tf, "  This strictly accounts for the 252 business days in a standard trading year.", 16, False, LIGHT)

    # -------------------------------------------------------------
    # SLIDE 5: Performance Scorecarding
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0.5, 0.5, 0.08, 0.6, *ACCENT)
    add_textbox(sl, 0.8, 0.4, 11, 0.7, "Mutual Fund Performance Scorecarding", 30, True, WHITE)
    
    tf = add_textbox(sl, 0.8, 1.5, 5.5, 5.0, "Ratios & Risk Metrics", 20, True, ACCENT)
    add_para(tf, "• Volatility (Annualised): Vol = Daily_Returns_std * sqrt(252).", 16, False, LIGHT)
    add_para(tf, "• Sharpe Ratio: Sharpe = (CAGR_3Yr - Rf) / Vol_Ann. Measures excess returns per unit of volatility (Rf = 6.0%).", 16, False, LIGHT)
    add_para(tf, "• Sortino Ratio: Sortino = (CAGR_3Yr - Rf) / Downside_Vol_Ann. Focuses only on negative volatility.", 16, False, LIGHT)
    add_para(tf, "• Alpha & Beta: Estimated via OLS linear regression of fund returns against the Nifty benchmark indices.", 16, False, LIGHT)
    add_para(tf, "• Maximum Drawdown: Measures historical peak-to-trough drop and tracks recovery dates.", 16, False, LIGHT)
    
    tf2 = add_textbox(sl, 6.8, 1.5, 5.5, 5.0, "Scorecard Rankings", 20, True, GOLD)
    add_para(tf2, "• Best Performers: SBI Bluechip and ICICI Bluechip display robust 3Yr CAGRs with high Sharpe ratios.", 16, False, LIGHT)
    add_para(tf2, "• Low-Cost Leadership: Direct plans consistently show superior CAGRs due to lower expense ratios (often 0.5% - 1.0% cheaper than regular plans).", 16, False, LIGHT)
    add_para(tf2, "• Scorecard CSV: Exported to data/processed/fund_scorecard.csv.", 16, False, LIGHT)

    # -------------------------------------------------------------
    # SLIDE 6: Advanced Financial Engineering
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0.5, 0.5, 0.08, 0.6, *ACCENT)
    add_textbox(sl, 0.8, 0.4, 11, 0.7, "Advanced Risk & Portfolio Engineering", 30, True, WHITE)
    
    tf = add_textbox(sl, 0.8, 1.5, 5.5, 5.0, "Value at Risk (VaR) & CVaR (D6)", 20, True, ACCENT)
    add_para(tf, "• 95% Daily VaR: The 5th percentile of daily returns, indicating the threshold loss that is exceeded only 5% of the time.", 16, False, LIGHT)
    add_para(tf, "• Conditional VaR (CVaR): The mean of the returns in the 5% tail, representing expected losses in extreme market scenarios.", 16, False, LIGHT)
    add_para(tf, "• Output: High beta small-cap funds show daily VaR losses exceeding 1.8%, while gilt debt funds remain below 0.3%.", 16, False, LIGHT)
    
    tf2 = add_textbox(sl, 6.8, 1.5, 5.5, 5.0, "HHI Concentration (D6)", 20, True, GOLD)
    add_para(tf2, "• Herfindahl-Hirschman Index: Sector concentration = sum(Sector_Weight_pct ^ 2).", 16, False, LIGHT)
    add_para(tf2, "• Score Ranges: Scores over 2,000 signify heavy sector focus. Scores below 1,000 show multi-sector diversification.", 16, False, LIGHT)
    add_para(tf2, "• Insights: Equity Large Cap funds showcase diversified sector holdings, while thematic schemes are heavily concentrated.", 16, False, LIGHT)

    # -------------------------------------------------------------
    # SLIDE 7: Projections & Markowitz Optimization
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0.5, 0.5, 0.08, 0.6, *ACCENT)
    add_textbox(sl, 0.8, 0.4, 11, 0.7, "NAV Projections & Portfolio Optimization", 30, True, WHITE)
    
    tf = add_textbox(sl, 0.8, 1.5, 5.5, 5.0, "Monte Carlo Simulation (B3)", 20, True, ACCENT)
    add_para(tf, "• Method: Geometric Brownian Motion (GBM) with historical drift (mu) and volatility (sigma) parameters.", 16, False, LIGHT)
    add_para(tf, "• Projection: Simulates 1,000 daily paths over a 5-year horizon.", 16, False, LIGHT)
    add_para(tf, "• Output: Generates median expected paths along with 90% confidence bands showing path uncertainty.", 16, False, LIGHT)
    
    tf2 = add_textbox(sl, 6.8, 1.5, 5.5, 5.0, "Efficient Frontier (B4)", 20, True, GOLD)
    add_para(tf2, "• Markowitz Frontier: Simulates weight combinations for 5 selected large-cap funds.", 16, False, LIGHT)
    add_para(tf2, "• Objective: Compute annualised expected returns and standard deviations for portfolios.", 16, False, LIGHT)
    add_para(tf2, "• Maximum Sharpe: Identifies optimal weight allocations (e.g. favoring SBI, ICICI, and Nippon Large Cap for best risk-adjusted performance).", 16, False, LIGHT)

    # -------------------------------------------------------------
    # SLIDE 8: Streamlit Dashboard Features
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0.5, 0.5, 0.08, 0.6, *ACCENT)
    add_textbox(sl, 0.8, 0.4, 11, 0.7, "Interactive Streamlit Web Dashboard", 30, True, WHITE)
    
    tf = add_textbox(sl, 0.8, 1.5, 11.5, 5.0, "High-End Features & Widgets", 20, True, ACCENT)
    add_para(tf, "• Multi-Page Layout: Organized into Executive Scorecard, NAV Visualizer, Sector Concentration, and Financial Engineering tabs.", 16, False, LIGHT)
    add_para(tf, "• Slicers (D5/B2): Every page contains at least 2 active input filters (e.g. fund house, category, scheme selector, date range, risk appetite).", 16, False, LIGHT)
    add_para(tf, "• Plots: Interactive Plotly line charts, horizontal bars, donuts, and treemaps reacting immediately to filter changes.", 16, False, LIGHT)
    add_para(tf, "• Portfolio Tools: Dynamic Monte Carlo projections slider and Markowitz frontier plots based on user-controlled parameters.", 16, False, LIGHT)
    add_para(tf, "• Launch Command: run using `streamlit run dashboard/app.py` for local testing.", 16, False, LIGHT)

    # -------------------------------------------------------------
    # SLIDE 9: Investor Behavior Analytics
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0.5, 0.5, 0.08, 0.6, *ACCENT)
    add_textbox(sl, 0.8, 0.4, 11, 0.7, "Investor Behavior Analytics", 30, True, WHITE)
    
    tf = add_textbox(sl, 0.8, 1.5, 5.5, 5.0, "Cohort Analysis (D6)", 20, True, ACCENT)
    add_para(tf, "• Segmentation: Segments investors based on their first transaction year (2024 vs 2025 cohorts).", 16, False, LIGHT)
    add_para(tf, "• Cohort behavior: Computes count, average SIP amount, total cumulative investments, and top fund house choice.", 16, False, LIGHT)
    add_para(tf, "• Summary: 2024 cohort displays higher cumulative volumes, while 2025 cohort demonstrates higher average monthly SIP values.", 16, False, LIGHT)
    
    tf2 = add_textbox(sl, 6.8, 1.5, 5.5, 5.0, "SIP Continuity & Churn (D6)", 20, True, GOLD)
    add_para(tf2, "• Criteria: Filters investors with 6 or more consecutive SIP cycles.", 16, False, LIGHT)
    add_para(tf2, "• Missed Payment Flag: Measures transaction date gaps. If the maximum gap between payments exceeds 35 days, the account is flagged as 'At Risk'.", 16, False, LIGHT)
    add_para(tf2, "• Output: Generates data/processed/sip_continuity.csv containing target accounts for churn-prevention emails.", 16, False, LIGHT)

    # -------------------------------------------------------------
    # SLIDE 10: Deliverables Verification Summary
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0.5, 0.5, 0.08, 0.6, *ACCENT)
    add_textbox(sl, 0.8, 0.4, 11, 0.7, "Deliverables Verification", 30, True, WHITE)
    
    deliverables = [
        ("scripts/etl_pipeline.py", "Fetches NAV, cleans holidays, loads SQLite DB"),
        ("data/db/bluestock_mf.db", "Ignored from git (*.db in .gitignore)"),
        ("sql/schema.sql & queries.sql", "SQLite database DDL schema and 10 analytical queries"),
        ("notebooks/ (01 to 05)", "Five clean, structured Jupyter notebooks covering ingestion, cleaning, EDA, performance, and advanced analytics"),
        ("dashboard/app.py", "Streamlit application containing Scorecard, NAV/Benchmark comparison, HHI charts, Monte Carlo, and Frontier"),
        ("scripts/email_report.py", "Automated weekly HTML performance report generator"),
        ("reports/Final_Report.pdf", "Detailed analytical summary report"),
    ]
    
    add_textbox(sl, 1.0, 1.6, 4.0, 0.4, "DELIVERABLE FILE", 16, True, ACCENT)
    add_textbox(sl, 5.5, 1.6, 7.0, 0.4, "VERIFIED STATUS / DESCRIPTION", 16, True, ACCENT)
    add_shape_bar(sl, 1.0, 2.0, 11.3, 0.02, *ACCENT)
    
    y = 2.15
    for path, desc in deliverables:
        add_textbox(sl, 1.0, y, 4.0, 0.35, path, 15, True, GOLD)
        add_textbox(sl, 5.5, y, 7.0, 0.35, desc, 15, False, LIGHT)
        y += 0.52
        
    add_textbox(sl, 1.0, y + 0.3, 11, 0.5, "Project fully completed and deployed to GitHub! ✓", 18, True, GREEN)

    # -------------------------------------------------------------
    # SLIDE 11: Thank You
    # -------------------------------------------------------------
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, *BG)
    add_shape_bar(sl, 0, 0, 13.333, 0.08, *ACCENT)
    add_shape_bar(sl, 0, 7.42, 13.333, 0.08, *ACCENT2)
    
    add_textbox(sl, 1, 2.5, 11, 1.0, "Thank You!", 44, True, ACCENT, PP_ALIGN.CENTER)
    add_textbox(sl, 1, 3.8, 11, 0.6, "Mutual Fund Capstone Platform Completed Successfully", 24, False, WHITE, PP_ALIGN.CENTER)
    add_textbox(sl, 1, 4.8, 11, 0.5, "Deployable on Streamlit & SQLite Core", 20, False, LIGHT, PP_ALIGN.CENTER)

    # Save
    output_path = REPORTS_DIR / "Presentation.pptx"
    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
